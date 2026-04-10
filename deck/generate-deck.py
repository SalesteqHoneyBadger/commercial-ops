#!/usr/bin/env python3
"""Generate the Salesteq event deck as .pptx"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BG = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
GREY = RGBColor(136, 136, 136)
DIM = RGBColor(68, 68, 68)
ORANGE = RGBColor(232, 114, 12)

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

def add_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return slide

def add_text(slide, text, left, top, width, height, size=32, color=WHITE, bold=False, align=PP_ALIGN.CENTER, font_name="Helvetica"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return tf

def add_para(tf, text, size=32, color=WHITE, bold=False, align=PP_ALIGN.CENTER):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Helvetica"
    p.alignment = align
    return p

# ── SLIDE 1: Title ──
s = add_slide()
add_text(s, "Salesteq", 0, 3.2, 16, 1.5, size=72, bold=True)
add_text(s, "AI Commercial Operations", 0, 5.0, 16, 0.8, size=28, color=GREY)

# ── SLIDE 2: Value prop ──
s = add_slide()
tf = add_text(s, "One platform replaces your", 2, 2.8, 12, 3.5, size=48, bold=True)
add_para(tf, "sales, marketing, and service teams", size=48, bold=True)
add_para(tf, "with AI agents.", size=48, color=ORANGE, bold=True)

# ── SLIDE 3: The problem ──
s = add_slide()
add_text(s, "European car dealers face", 0, 1.5, 16, 0.8, size=28, color=GREY)

stats = [("4", "Languages in\nSwitzerland"), ("200+", "Locations\nper group"), ("24/7", "Customer\nexpectations"), ("0", "Staff\nto hire")]
for i, (num, label) in enumerate(stats):
    x = 2.0 + i * 3.2
    c = DIM if num == "0" else WHITE
    add_text(s, num, x, 3.2, 2.5, 1.2, size=56, bold=True, color=c, font_name="Menlo")
    add_text(s, label, x, 4.8, 2.5, 1.0, size=14, color=GREY)

# ── SLIDE 4: The platform ──
s = add_slide()
add_text(s, "The Platform", 0, 1.2, 16, 1, size=48, bold=True)

products = [
    ("Badger", "AI agents that plan, act, and learn.\nThey execute sales, marketing, and\nservice tasks autonomously."),
    ("Index", "Commercial intelligence engine.\n6.8M companies, 14M contacts.\nSub-100ms search."),
    ("Kol", "Voice and chat in any language.\nCalls, WhatsApp, web chat.\nNative quality."),
    ("Yeda", "Self-building knowledge base.\nAgents learn from every customer\ninteraction automatically."),
]
for i, (name, desc) in enumerate(products):
    x = 1.5 + i * 3.5
    add_text(s, name, x, 3.5, 3, 0.6, size=22, color=ORANGE, bold=True, align=PP_ALIGN.LEFT)
    add_text(s, desc, x, 4.2, 3, 1.5, size=14, color=GREY, align=PP_ALIGN.LEFT)

# ── SLIDE 5: NAGHI ──
s = add_slide()
add_text(s, "Proof", 0, 1.5, 16, 0.6, size=28, color=GREY)
add_text(s, "NAGHI Motors", 0, 2.2, 16, 1, size=48, bold=True)

case_stats = [("11+", "Brands"), ("250K", "Vehicles / Year"), ("2", "Languages")]
for i, (num, label) in enumerate(case_stats):
    x = 3.5 + i * 3.2
    add_text(s, num, x, 3.8, 2.5, 1, size=64, bold=True, color=ORANGE, font_name="Menlo")
    add_text(s, label, x, 5.2, 2.5, 0.5, size=13, color=GREY)

add_text(s, "Exclusive BMW dealer in Saudi Arabia. Salesteq AI handles their entire", 3, 6.2, 10, 1.2, size=18, color=GREY)
tf = add_text(s, "customer journey across Arabic and English, 24/7.", 3, 6.7, 10, 0.6, size=18, color=GREY)

# ── SLIDE 6: What you're about to see ──
s = add_slide()
add_text(s, "What you're about to see", 0, 2.2, 16, 1, size=48, bold=True)
add_text(s, "5 AI agents will launch a full commercial campaign\ntargeting European automotive dealers.\nResearch, outreach, quality review, landing page, PDF proposals.\nAll in 15 minutes.", 3, 3.8, 10, 2, size=20, color=GREY)
tf = add_text(s, "The human equivalent:", 0, 6.2, 10, 0.6, size=28, color=GREY)
add_text(s, "60+ hours", 10, 6.2, 6, 0.6, size=28, bold=True)

# ── SLIDE 7: Live Demo ──
s = add_slide()
add_text(s, "Live", 0, 2.5, 16, 0.8, size=28, color=GREY)
add_text(s, "Demo", 0, 3.5, 16, 1.5, size=72, bold=True)
add_text(s, "automotive.salesteq.com", 0, 5.5, 16, 0.8, size=24, color=ORANGE, font_name="Menlo")

# Save
output = "/root/commercial-ops/deck/Salesteq-Demo.pptx"
prs.save(output)
print(f"Generated: {output}")
